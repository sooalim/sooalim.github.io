"""
Data Extraction Module for RAG Agent

This module handles the extraction of text content from various document formats
including PDF, Word, Excel, PowerPoint, HTML, and plain text files.
"""

import asyncio
import logging
import os
import tempfile
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
from urllib.parse import urlparse

import aiofiles
import chardet
import magic
import pandas as pd
import requests
from azure.core.exceptions import AzureError
from azure.identity import DefaultAzureCredential
from azure.keyvault.secrets import SecretClient
from azure.storage.blob import BlobServiceClient
from bs4 import BeautifulSoup
from docx import Document
from html2text import html2text
from PyPDF2 import PdfReader
from pptx import Presentation

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


@dataclass
class DocumentMetadata:
    """Metadata for extracted documents."""
    source_path: str
    file_type: str
    file_size: int
    title: Optional[str] = None
    author: Optional[str] = None
    created_date: Optional[datetime] = None
    modified_date: Optional[datetime] = None
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    language: Optional[str] = None
    encoding: Optional[str] = None
    extraction_date: datetime = field(default_factory=datetime.now)
    custom_metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class ExtractedDocument:
    """Container for extracted document content and metadata."""
    content: str
    metadata: DocumentMetadata
    sections: List[Dict[str, Any]] = field(default_factory=list)
    tables: List[Dict[str, Any]] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)
    links: List[Dict[str, Any]] = field(default_factory=list)


class DocumentExtractor(ABC):
    """Abstract base class for document extractors."""
    
    @abstractmethod
    def can_extract(self, file_path: str, mime_type: str) -> bool:
        """Check if this extractor can handle the given file."""
        pass
    
    @abstractmethod
    async def extract(self, file_path: str, **kwargs) -> ExtractedDocument:
        """Extract content from the document."""
        pass


class PDFExtractor(DocumentExtractor):
    """Extractor for PDF documents."""
    
    def can_extract(self, file_path: str, mime_type: str) -> bool:
        """Check if this extractor can handle PDF files."""
        return mime_type == 'application/pdf' or file_path.lower().endswith('.pdf')
    
    async def extract(self, file_path: str, **kwargs) -> ExtractedDocument:
        """Extract text content from PDF files."""
        try:
            logger.info(f"Extracting PDF: {file_path}")
            
            # Read PDF content
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                
                # Extract text from all pages
                text_content = []
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(page_text)
                
                content = '\n\n'.join(text_content)
                
                # Extract metadata
                pdf_info = pdf_reader.metadata
                file_stat = os.stat(file_path)
                
                metadata = DocumentMetadata(
                    source_path=file_path,
                    file_type='pdf',
                    file_size=file_stat.st_size,
                    title=pdf_info.get('/Title'),
                    author=pdf_info.get('/Author'),
                    created_date=self._parse_pdf_date(pdf_info.get('/CreationDate')),
                    modified_date=self._parse_pdf_date(pdf_info.get('/ModDate')),
                    page_count=len(pdf_reader.pages),
                    word_count=len(content.split()) if content else 0
                )
                
                return ExtractedDocument(content=content, metadata=metadata)
                
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {str(e)}")
            raise
    
    def _parse_pdf_date(self, date_str: Optional[str]) -> Optional[datetime]:
        """Parse PDF date format to datetime."""
        if not date_str:
            return None
        
        try:
            # PDF date format: D:YYYYMMDDHHmmSSOHH'mm'
            if date_str.startswith('D:'):
                date_str = date_str[2:]
            
            # Parse basic date part
            if len(date_str) >= 8:
                return datetime.strptime(date_str[:8], '%Y%m%d')
        except Exception as e:
            logger.warning(f"Failed to parse PDF date {date_str}: {str(e)}")
        
        return None


class WordExtractor(DocumentExtractor):
    """Extractor for Word documents."""
    
    def can_extract(self, file_path: str, mime_type: str) -> bool:
        """Check if this extractor can handle Word files."""
        return (mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                              'application/msword'] or
                file_path.lower().endswith(('.docx', '.doc')))
    
    async def extract(self, file_path: str, **kwargs) -> ExtractedDocument:
        """Extract text content from Word documents."""
        try:
            logger.info(f"Extracting Word document: {file_path}")
            
            # Read Word document
            doc = Document(file_path)
            
            # Extract text content
            paragraphs = []
            tables = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    paragraphs.append(text)
            
            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                
                if table_data:
                    tables.append({
                        'data': table_data,
                        'rows': len(table_data),
                        'columns': len(table_data[0]) if table_data else 0
                    })
            
            content = '\n\n'.join(paragraphs)
            
            # Extract metadata
            file_stat = os.stat(file_path)
            core_props = doc.core_properties
            
            metadata = DocumentMetadata(
                source_path=file_path,
                file_type='docx',
                file_size=file_stat.st_size,
                title=core_props.title,
                author=core_props.author,
                created_date=core_props.created,
                modified_date=core_props.modified,
                word_count=len(content.split()) if content else 0
            )
            
            return ExtractedDocument(
                content=content,
                metadata=metadata,
                tables=tables
            )
            
        except Exception as e:
            logger.error(f"Error extracting Word document {file_path}: {str(e)}")
            raise


class ExcelExtractor(DocumentExtractor):
    """Extractor for Excel documents."""
    
    def can_extract(self, file_path: str, mime_type: str) -> bool:
        """Check if this extractor can handle Excel files."""
        return (mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                              'application/vnd.ms-excel'] or
                file_path.lower().endswith(('.xlsx', '.xls')))
    
    async def extract(self, file_path: str, **kwargs) -> ExtractedDocument:
        """Extract content from Excel files."""
        try:
            logger.info(f"Extracting Excel document: {file_path}")
            
            # Read Excel file
            excel_file = pd.ExcelFile(file_path)
            
            content_parts = []
            tables = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert DataFrame to text
                sheet_content = f"Sheet: {sheet_name}\n"
                sheet_content += df.to_string(index=False, na_rep='')
                content_parts.append(sheet_content)
                
                # Store table data
                tables.append({
                    'sheet_name': sheet_name,
                    'data': df.values.tolist(),
                    'columns': df.columns.tolist(),
                    'rows': len(df),
                    'columns_count': len(df.columns)
                })
            
            content = '\n\n'.join(content_parts)
            
            # Extract metadata
            file_stat = os.stat(file_path)
            
            metadata = DocumentMetadata(
                source_path=file_path,
                file_type='xlsx',
                file_size=file_stat.st_size,
                word_count=len(content.split()) if content else 0,
                custom_metadata={'sheet_count': len(excel_file.sheet_names)}
            )
            
            return ExtractedDocument(
                content=content,
                metadata=metadata,
                tables=tables
            )
            
        except Exception as e:
            logger.error(f"Error extracting Excel document {file_path}: {str(e)}")
            raise


class PowerPointExtractor(DocumentExtractor):
    """Extractor for PowerPoint presentations."""
    
    def can_extract(self, file_path: str, mime_type: str) -> bool:
        """Check if this extractor can handle PowerPoint files."""
        return (mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                              'application/vnd.ms-powerpoint'] or
                file_path.lower().endswith(('.pptx', '.ppt')))
    
    async def extract(self, file_path: str, **kwargs) -> ExtractedDocument:
        """Extract content from PowerPoint presentations."""
        try:
            logger.info(f"Extracting PowerPoint presentation: {file_path}")
            
            # Read PowerPoint file
            prs = Presentation(file_path)
            
            slides_content = []
            sections = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                slide_content = '\n'.join(slide_text)
                if slide_content:
                    slides_content.append(f"Slide {slide_num}:\n{slide_content}")
                    sections.append({
                        'type': 'slide',
                        'number': slide_num,
                        'content': slide_content
                    })
            
            content = '\n\n'.join(slides_content)
            
            # Extract metadata
            file_stat = os.stat(file_path)
            
            metadata = DocumentMetadata(
                source_path=file_path,
                file_type='pptx',
                file_size=file_stat.st_size,
                word_count=len(content.split()) if content else 0,
                custom_metadata={'slide_count': len(prs.slides)}
            )
            
            return ExtractedDocument(
                content=content,
                metadata=metadata,
                sections=sections
            )
            
        except Exception as e:
            logger.error(f"Error extracting PowerPoint presentation {file_path}: {str(e)}")
            raise


class HTMLExtractor(DocumentExtractor):
    """Extractor for HTML documents."""
    
    def can_extract(self, file_path: str, mime_type: str) -> bool:
        """Check if this extractor can handle HTML files."""
        return (mime_type in ['text/html', 'application/xhtml+xml'] or
                file_path.lower().endswith(('.html', '.htm', '.xhtml')))
    
    async def extract(self, file_path: str, **kwargs) -> ExtractedDocument:
        """Extract content from HTML files."""
        try:
            logger.info(f"Extracting HTML document: {file_path}")
            
            # Read HTML file
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                html_content = await file.read()
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract metadata
            title = soup.find('title')
            title_text = title.get_text().strip() if title else None
            
            # Extract links
            links = []
            for link in soup.find_all('a', href=True):
                links.append({
                    'text': link.get_text().strip(),
                    'url': link['href']
                })
            
            # Extract tables
            tables = []
            for table in soup.find_all('table'):
                table_data = []
                for row in table.find_all('tr'):
                    row_data = []
                    for cell in row.find_all(['td', 'th']):
                        row_data.append(cell.get_text().strip())
                    if row_data:
                        table_data.append(row_data)
                
                if table_data:
                    tables.append({
                        'data': table_data,
                        'rows': len(table_data),
                        'columns': len(table_data[0]) if table_data else 0
                    })
            
            # Convert to plain text
            content = html2text(html_content)
            
            # Extract metadata
            file_stat = os.stat(file_path)
            
            metadata = DocumentMetadata(
                source_path=file_path,
                file_type='html',
                file_size=file_stat.st_size,
                title=title_text,
                word_count=len(content.split()) if content else 0
            )
            
            return ExtractedDocument(
                content=content,
                metadata=metadata,
                tables=tables,
                links=links
            )
            
        except Exception as e:
            logger.error(f"Error extracting HTML document {file_path}: {str(e)}")
            raise


class TextExtractor(DocumentExtractor):
    """Extractor for plain text documents."""
    
    def can_extract(self, file_path: str, mime_type: str) -> bool:
        """Check if this extractor can handle text files."""
        return (mime_type.startswith('text/') or
                file_path.lower().endswith(('.txt', '.md', '.csv', '.json', '.xml', '.log')))
    
    async def extract(self, file_path: str, **kwargs) -> ExtractedDocument:
        """Extract content from text files."""
        try:
            logger.info(f"Extracting text document: {file_path}")
            
            # Detect encoding
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                encoding_result = chardet.detect(raw_data)
                encoding = encoding_result.get('encoding', 'utf-8')
            
            # Read text content
            async with aiofiles.open(file_path, 'r', encoding=encoding) as file:
                content = await file.read()
            
            # Extract metadata
            file_stat = os.stat(file_path)
            
            metadata = DocumentMetadata(
                source_path=file_path,
                file_type='text',
                file_size=file_stat.st_size,
                word_count=len(content.split()) if content else 0,
                encoding=encoding
            )
            
            return ExtractedDocument(content=content, metadata=metadata)
            
        except Exception as e:
            logger.error(f"Error extracting text document {file_path}: {str(e)}")
            raise


class DataExtractionManager:
    """Main class for managing document extraction operations."""
    
    def __init__(self, key_vault_url: Optional[str] = None):
        """Initialize the data extraction manager."""
        self.key_vault_url = key_vault_url
        self.credential = DefaultAzureCredential()
        self.blob_client = None
        self.secret_client = None
        
        # Initialize extractors
        self.extractors = [
            PDFExtractor(),
            WordExtractor(),
            ExcelExtractor(),
            PowerPointExtractor(),
            HTMLExtractor(),
            TextExtractor()
        ]
        
        # Initialize Azure clients if key vault is provided
        if self.key_vault_url:
            self._initialize_azure_clients()
    
    def _initialize_azure_clients(self):
        """Initialize Azure service clients."""
        try:
            self.secret_client = SecretClient(
                vault_url=self.key_vault_url,
                credential=self.credential
            )
            
            # Get storage connection string from Key Vault
            storage_secret = self.secret_client.get_secret("storage-connection-string")
            self.blob_client = BlobServiceClient.from_connection_string(
                storage_secret.value
            )
            
            logger.info("Azure clients initialized successfully")
            
        except Exception as e:
            logger.error(f"Failed to initialize Azure clients: {str(e)}")
            raise
    
    def _get_mime_type(self, file_path: str) -> str:
        """Get MIME type of a file."""
        try:
            return magic.from_file(file_path, mime=True)
        except Exception as e:
            logger.warning(f"Failed to detect MIME type for {file_path}: {str(e)}")
            return 'application/octet-stream'
    
    def _get_extractor(self, file_path: str, mime_type: str) -> Optional[DocumentExtractor]:
        """Get appropriate extractor for a file."""
        for extractor in self.extractors:
            if extractor.can_extract(file_path, mime_type):
                return extractor
        return None
    
    async def extract_from_file(self, file_path: str, **kwargs) -> ExtractedDocument:
        """Extract content from a single file."""
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            
            mime_type = self._get_mime_type(file_path)
            extractor = self._get_extractor(file_path, mime_type)
            
            if not extractor:
                raise ValueError(f"No extractor available for file type: {mime_type}")
            
            logger.info(f"Using {extractor.__class__.__name__} for {file_path}")
            return await extractor.extract(file_path, **kwargs)
            
        except Exception as e:
            logger.error(f"Error extracting from file {file_path}: {str(e)}")
            raise
    
    async def extract_from_url(self, url: str, **kwargs) -> ExtractedDocument:
        """Extract content from a URL."""
        try:
            logger.info(f"Downloading content from URL: {url}")
            
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=self._get_suffix_from_url(url)) as temp_file:
                temp_file.write(response.content)
                temp_file_path = temp_file.name
            
            try:
                # Extract from temporary file
                extracted_doc = await self.extract_from_file(temp_file_path, **kwargs)
                
                # Update metadata with URL
                extracted_doc.metadata.source_path = url
                extracted_doc.metadata.custom_metadata['original_url'] = url
                
                return extracted_doc
                
            finally:
                # Clean up temporary file
                os.unlink(temp_file_path)
                
        except Exception as e:
            logger.error(f"Error extracting from URL {url}: {str(e)}")
            raise
    
    def _get_suffix_from_url(self, url: str) -> str:
        """Get file suffix from URL."""
        parsed_url = urlparse(url)
        path = parsed_url.path
        
        if '.' in path:
            return '.' + path.split('.')[-1]
        return '.html'  # Default to HTML
    
    async def extract_from_blob(self, container_name: str, blob_name: str, **kwargs) -> ExtractedDocument:
        """Extract content from Azure Blob Storage."""
        try:
            if not self.blob_client:
                raise ValueError("Blob client not initialized. Please provide key_vault_url.")
            
            logger.info(f"Downloading blob: {container_name}/{blob_name}")
            
            # Download blob to temporary file
            blob_client = self.blob_client.get_blob_client(
                container=container_name,
                blob=blob_name
            )
            
            with tempfile.NamedTemporaryFile(delete=False, suffix=self._get_suffix_from_blob(blob_name)) as temp_file:
                blob_data = blob_client.download_blob()
                temp_file.write(blob_data.readall())
                temp_file_path = temp_file.name
            
            try:
                # Extract from temporary file
                extracted_doc = await self.extract_from_file(temp_file_path, **kwargs)
                
                # Update metadata with blob info
                extracted_doc.metadata.source_path = f"blob://{container_name}/{blob_name}"
                extracted_doc.metadata.custom_metadata['container_name'] = container_name
                extracted_doc.metadata.custom_metadata['blob_name'] = blob_name
                
                return extracted_doc
                
            finally:
                # Clean up temporary file
                os.unlink(temp_file_path)
                
        except Exception as e:
            logger.error(f"Error extracting from blob {container_name}/{blob_name}: {str(e)}")
            raise
    
    def _get_suffix_from_blob(self, blob_name: str) -> str:
        """Get file suffix from blob name."""
        if '.' in blob_name:
            return '.' + blob_name.split('.')[-1]
        return '.bin'
    
    async def extract_batch(self, file_paths: List[str], **kwargs) -> List[ExtractedDocument]:
        """Extract content from multiple files."""
        tasks = []
        for file_path in file_paths:
            task = self.extract_from_file(file_path, **kwargs)
            tasks.append(task)
        
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        # Filter out exceptions and log them
        extracted_docs = []
        for i, result in enumerate(results):
            if isinstance(result, Exception):
                logger.error(f"Failed to extract from {file_paths[i]}: {str(result)}")
            else:
                extracted_docs.append(result)
        
        return extracted_docs
    
    async def extract_from_directory(self, directory_path: str, recursive: bool = True, **kwargs) -> List[ExtractedDocument]:
        """Extract content from all supported files in a directory."""
        try:
            if not os.path.exists(directory_path):
                raise FileNotFoundError(f"Directory not found: {directory_path}")
            
            file_paths = []
            
            if recursive:
                for root, dirs, files in os.walk(directory_path):
                    for file in files:
                        file_path = os.path.join(root, file)
                        file_paths.append(file_path)
            else:
                for item in os.listdir(directory_path):
                    item_path = os.path.join(directory_path, item)
                    if os.path.isfile(item_path):
                        file_paths.append(item_path)
            
            # Filter files that have supported extractors
            supported_files = []
            for file_path in file_paths:
                mime_type = self._get_mime_type(file_path)
                if self._get_extractor(file_path, mime_type):
                    supported_files.append(file_path)
            
            logger.info(f"Found {len(supported_files)} supported files in {directory_path}")
            return await self.extract_batch(supported_files, **kwargs)
            
        except Exception as e:
            logger.error(f"Error extracting from directory {directory_path}: {str(e)}")
            raise


# Example usage and testing
async def main():
    """Example usage of the data extraction module."""
    
    # Initialize extraction manager
    manager = DataExtractionManager()
    
    # Example: Extract from a single file
    try:
        doc = await manager.extract_from_file("sample.pdf")
        print(f"Extracted {len(doc.content)} characters from {doc.metadata.source_path}")
        print(f"File type: {doc.metadata.file_type}")
        print(f"Page count: {doc.metadata.page_count}")
        
    except Exception as e:
        print(f"Error: {str(e)}")
    
    # Example: Extract from directory
    try:
        docs = await manager.extract_from_directory("./documents", recursive=True)
        print(f"Extracted content from {len(docs)} documents")
        
        for doc in docs:
            print(f"- {doc.metadata.source_path}: {len(doc.content)} characters")
            
    except Exception as e:
        print(f"Error: {str(e)}")


if __name__ == "__main__":
    asyncio.run(main())
